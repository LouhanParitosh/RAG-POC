"""
RAG POC — PowerPoint presentation generator.
Run: python create_ppt.py
Output: RAG_POC_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
CARD        = RGBColor(0x1E, 0x2D, 0x45)   # dark card
BLUE        = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
GREEN       = RGBColor(0x34, 0xD3, 0x99)   # emerald
YELLOW      = RGBColor(0xFB, 0xBF, 0x24)   # amber
PINK        = RGBColor(0xF4, 0x72, 0xB6)   # pink
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)   # violet
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0xCB, 0xD5, 0xE1)

W = Inches(13.33)
H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def text(slide, txt, l, t, w, h,
         size=14, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def multiline(slide, lines, l, t, w, h,
              size=12, color=GREY, bold_first=False, spacing=4):
    """Render a list of strings as stacked paragraphs in one textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)


# ── Composite helpers ──────────────────────────────────────────────────────────

def top_bar(slide, color=BLUE):
    rect(slide, 0, 0, 13.33, 0.06, color)


def slide_heading(slide, title, subtitle=None, bar_color=BLUE):
    top_bar(slide, bar_color)
    text(slide, title, 0.5, 0.25, 12.3, 0.9, size=32, bold=True, color=WHITE)
    if subtitle:
        text(slide, subtitle, 0.5, 1.1, 12.3, 0.5, size=15, color=GREY)
    rect(slide, 0.5, 1.65, 0.9, 0.05, bar_color)


def card(slide, l, t, w, h, title, lines,
         title_color=BLUE, body_color=GREY, bg_color=CARD, body_size=11):
    rect(slide, l, t, w, h, bg_color)
    text(slide, title, l + 0.15, t + 0.12, w - 0.25, 0.38,
         size=13, bold=True, color=title_color)
    body_top = t + 0.52
    body_h = h - 0.6
    multiline(slide, lines, l + 0.15, body_top, w - 0.25, body_h,
              size=body_size, color=body_color, spacing=3)


def row_of_cards(slide, top, items, card_h=1.8, gap=0.15):
    """items = list of (title, lines, title_color).  Auto-sizes widths."""
    n = len(items)
    total_gap = gap * (n - 1)
    card_w = (13.33 - 0.3 * 2 - total_gap) / n
    for i, (title, lines, tc) in enumerate(items):
        l = 0.3 + i * (card_w + gap)
        card(slide, l, top, card_w, card_h, title, lines, title_color=tc)


def code_box(slide, code_lines, l, t, w, h, label=None):
    rect(slide, l, t, w, h, CARD)
    if label:
        text(slide, label, l + 0.12, t + 0.08, w - 0.2, 0.3,
             size=10, bold=True, color=GREY)
        multiline(slide, code_lines, l + 0.12, t + 0.38, w - 0.2, h - 0.45,
                  size=10, color=YELLOW, spacing=2)
    else:
        multiline(slide, code_lines, l + 0.12, t + 0.12, w - 0.2, h - 0.2,
                  size=10, color=YELLOW, spacing=2)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── 1. TITLE ───────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
top_bar(s, BLUE)
rect(s, 0, 6.55, 13.33, 0.95, CARD)
text(s, "RAG", 1.0, 0.9, 12, 2.6,
     size=110, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
text(s, "Retrieval-Augmented Generation", 1.0, 3.45, 12, 0.8,
     size=32, color=WHITE, align=PP_ALIGN.LEFT)
text(s, "A complete technical deep-dive — concepts, models, algorithms & architecture",
     1.0, 4.2, 12, 0.6, size=17, color=GREY, align=PP_ALIGN.LEFT)
text(s, "Embeddings  ·  Vector Search  ·  HNSW  ·  LLMs  ·  Chunking  ·  Guardrails  ·  RAG Variants",
     1.0, 6.65, 12, 0.45, size=13, color=GREEN, align=PP_ALIGN.LEFT)


# ── 2. WHAT IS RAG ─────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "What is RAG?",
              "Teaching an LLM to answer from YOUR documents instead of its training data")

text(s, "The core problem with plain LLMs:", 0.5, 1.85, 12, 0.4,
     size=14, bold=True, color=YELLOW)
text(s,
     "A standard LLM like GPT-4 or llama3.2 only knows what it was trained on — data from months or "
     "years ago. It cannot access your private PDFs, internal reports, or real-time data. "
     "Worse, it will confidently make up answers (hallucinate) when it doesn't know something.",
     0.5, 2.25, 12.3, 0.75, size=13, color=GREY)

row_of_cards(s, 3.1, [
    ("① RETRIEVE",
     ["Convert the user's question to a vector (embedding).",
      "Search the vector database for the most",
      "similar document chunks.",
      "Return Top-K chunks ranked by similarity score."],
     BLUE),
    ("② AUGMENT",
     ["Inject the retrieved chunks into the LLM prompt",
      "alongside the original question.",
      "The LLM now has the exact context it needs.",
      "This is the 'augmented prompt'."],
     GREEN),
    ("③ GENERATE",
     ["The LLM reads the context + question.",
      "Produces a grounded, accurate answer.",
      "Cites source documents.",
      "Refuses to answer if context is insufficient."],
     YELLOW),
], card_h=2.1)

text(s, "RAG = best of both worlds: LLM's language skill + your document's up-to-date facts.",
     0.5, 5.3, 12.3, 0.45, size=13, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

row_of_cards(s, 5.85, [
    ("Without RAG",
     ["Uses training data only",
      "Knowledge cutoff date",
      "May hallucinate",
      "Cannot use private data"],
     PINK),
    ("With RAG",
     ["Uses your documents",
      "Always up-to-date (re-ingest anytime)",
      "Grounded — cites sources",
      "Works with private/confidential data"],
     GREEN),
    ("Use cases",
     ["Customer support bots over documentation",
      "Legal document Q&A",
      "Medical literature search",
      "Enterprise knowledge bases"],
     BLUE),
    ("Limitations",
     ["Quality depends on chunking strategy",
      "Embedding model determines search accuracy",
      "LLM context window still limited",
      "Scanned PDFs need OCR first"],
     YELLOW),
], card_h=1.5)


# ── 3. ARCHITECTURE ────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
top_bar(s, GREEN)
text(s, "RAG Architecture", 0.5, 0.15, 12, 0.7,
     size=30, bold=True, color=WHITE)

text(s, "INGESTION PIPELINE  (runs once per document, or when files change)",
     0.3, 1.0, 12.7, 0.35, size=12, bold=True, color=BLUE)

ingestion_boxes = [
    ("Document\n(PDF/TXT/MD)", 0.3),
    ("Text\nChunker",          2.6),
    ("Embedding\nModel",       4.9),
    ("Vector\nDatabase",       7.2),
    ("Saved to\nDisk",         9.5),
]
for label, lpos in ingestion_boxes:
    rect(s, lpos, 1.4, 2.1, 0.85, CARD)
    text(s, label, lpos + 0.08, 1.42, 1.95, 0.8,
         size=11, color=BLUE, align=PP_ALIGN.CENTER)
for x in [2.45, 4.75, 7.05, 9.35]:
    text(s, "→", x, 1.58, 0.22, 0.5,
         size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

text(s, "QUERY PIPELINE  (runs on every question)",
     0.3, 2.45, 12.7, 0.35, size=12, bold=True, color=YELLOW)

query_boxes = [
    ("User\nQuestion",        0.3),
    ("Embed\nQuery",          2.3),
    ("Similarity\nSearch",    4.3),
    ("Top-K\nChunks",         6.3),
    ("Augmented\nPrompt",     8.3),
    ("LLM →\nAnswer",        10.5),
]
for label, lpos in query_boxes:
    rect(s, lpos, 2.85, 1.85, 0.85, CARD)
    text(s, label, lpos + 0.05, 2.87, 1.75, 0.8,
         size=11, color=YELLOW, align=PP_ALIGN.CENTER)
for x in [2.12, 4.12, 6.12, 8.12, 10.32]:
    text(s, "→", x, 3.03, 0.22, 0.5,
         size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

rect(s, 0.3, 3.85, 12.7, 0.04, BLUE)

row_of_cards(s, 4.0, [
    ("Critical rule",
     ["The SAME embedding model must be used for",
      "both ingestion and query. If you switch models,",
      "you MUST re-ingest all documents from scratch.",
      "Mixed embeddings produce garbage search results."],
     PINK),
    ("Why chunk instead of full docs?",
     ["LLMs have a context window limit (~4k–128k tokens).",
      "Smaller chunks = more precise similarity search.",
      "A relevant paragraph scores better than a whole chapter.",
      "Reduces cost and latency per query."],
     GREEN),
    ("What the LLM receives",
     ["NOT the full document.",
      "Only the Top-K most relevant chunks.",
      "Formatted as context blocks with source labels.",
      "Plus the original question."],
     BLUE),
    ("Persistent storage",
     ["Embeddings are computed once and stored.",
      "No re-embedding on each question.",
      "Only re-embed when documents change.",
      "Vector DB lives on disk between sessions."],
     YELLOW),
], card_h=2.9)


# ── 4. EMBEDDINGS ──────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Embeddings",
              "Converting text into numbers that capture semantic meaning")

text(s, "What is an embedding?", 0.5, 1.85, 12, 0.35,
     size=15, bold=True, color=BLUE)
text(s,
     "An embedding is a fixed-length list of decimal numbers (a vector) that encodes the *meaning* "
     "of a piece of text. The model is trained so that semantically similar texts produce vectors "
     "that are numerically close to each other — regardless of exact wording.",
     0.5, 2.2, 12.3, 0.7, size=13, color=GREY)

rect(s, 0.5, 3.0, 12.3, 0.55, CARD)
text(s, '"The patient has a high fever."  →  [0.23, -0.41, 0.87, 0.02, ... 768 numbers total]',
     0.65, 3.07, 12.0, 0.4, size=12, color=YELLOW)

rect(s, 0.5, 3.62, 12.3, 0.55, CARD)
text(s, '"The individual is experiencing elevated body temperature."  →  [0.24, -0.39, 0.85, 0.03, ...]  ← nearly identical!',
     0.65, 3.69, 12.0, 0.4, size=12, color=GREEN)

text(s, "Same meaning → similar numbers → found by same search   (even though not a single shared word)",
     0.5, 4.25, 12.3, 0.4, size=12, italic=True, color=GREY)

row_of_cards(s, 4.75, [
    ("Dimensions",
     ["Most modern models: 384–4096 dims.",
      "nomic-embed-text: 768 dims.",
      "text-embedding-3-large: 3072 dims.",
      "More dims ≠ always better.",
      "Tradeoff: quality vs storage vs speed."],
     BLUE),
    ("How they're trained",
     ["Trained on millions of sentence pairs",
      "labelled as similar or dissimilar.",
      "Contrastive learning: similar pairs pushed",
      "together, dissimilar pairs pushed apart.",
      "Result: geometry encodes meaning."],
     GREEN),
    ("What they enable",
     ["Semantic search (meaning, not keywords)",
      "Clustering similar documents",
      "Anomaly detection",
      "Recommendation systems",
      "Multilingual search"],
     YELLOW),
    ("Storage cost",
     ["1 chunk = 768 float32 numbers = ~3 KB.",
      "10,000 chunks ≈ 30 MB of vectors.",
      "1 million chunks ≈ 3 GB.",
      "ChromaDB stores vectors + text + metadata",
      "in a single SQLite file on disk."],
     PINK),
], card_h=2.1)


# ── 5. EMBEDDING MODELS ────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Types of Embedding Models",
              "Choosing the right model based on accuracy, cost, privacy, and speed",
              bar_color=GREEN)

models = [
    ("OpenAI",
     "text-embedding-3-small\ntext-embedding-3-large\ntext-embedding-ada-002",
     ["Best-in-class quality", "3-small: 1536 dims, 3-large: 3072 dims",
      "Paid API — ~$0.02 per million tokens",
      "Data sent to OpenAI servers",
      "Not suitable for private/sensitive data"],
     BLUE),
    ("Cohere",
     "embed-english-v3.0\nembed-multilingual-v3.0",
     ["Strong multilingual support (100+ languages)",
      "1024 dimensions",
      "Paid API with free tier",
      "Good for global / multilingual RAG",
      "Data sent to Cohere servers"],
     GREEN),
    ("Google",
     "text-embedding-004\ngemini-embedding-exp",
     ["High quality, large context window",
      "768 or 3072 dimensions",
      "Part of Google AI / Vertex AI",
      "Paid API — GCP pricing",
      "Data processed on Google infrastructure"],
     YELLOW),
    ("HuggingFace\n(Open Source)",
     "BAAI/bge-large-en-v1.5\nsentence-transformers\nE5-large-v2",
     ["Free — runs locally or on your server",
      "Wide variety: 384 to 1024 dims",
      "bge-large: top open-source benchmark scores",
      "Run via sentence-transformers library",
      "Full data privacy"],
     PURPLE),
    ("Ollama\n(Local)",
     "nomic-embed-text\nmxbai-embed-large\nsnowflake-arctic-embed",
     ["Completely local — no API, no cost",
      "nomic-embed-text: 768 dims, 274 MB",
      "mxbai-embed-large: 1024 dims, best local quality",
      "Zero data leaves your machine",
      "Best for private/offline RAG"],
     PINK),
]

for i, (provider, model_names, bullets, color) in enumerate(models):
    lpos = 0.3 + i * 2.57
    rect(s, lpos, 1.85, 2.42, 5.35, CARD)
    text(s, provider, lpos + 0.12, 1.95, 2.2, 0.5,
         size=13, bold=True, color=color)
    text(s, model_names, lpos + 0.12, 2.42, 2.2, 0.7,
         size=9, italic=True, color=GREY)
    rect(s, lpos + 0.12, 3.15, 2.18, 0.03, color)
    multiline(s, bullets, lpos + 0.12, 3.25, 2.2, 3.85,
              size=10, color=GREY, spacing=4)

text(s,
     "Rule: Use the SAME embedding model for both ingestion and query. "
     "Switching models requires re-ingesting all documents.",
     0.3, 7.2, 12.7, 0.35, size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── 6. VECTOR SEARCH & COSINE SIMILARITY ──────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Vector Search & Cosine Similarity",
              "How the system finds the most relevant chunks for any question",
              bar_color=PURPLE)

text(s, "Cosine Similarity — measuring the angle between two vectors", 0.5, 1.85, 12, 0.38,
     size=15, bold=True, color=PURPLE)
text(s,
     "Cosine similarity measures the angle between two vectors. "
     "A small angle (vectors pointing the same direction) = high similarity. "
     "A large angle = low similarity. The result is always between 0 and 1.",
     0.5, 2.23, 8.5, 0.75, size=13, color=GREY)

rect(s, 9.2, 1.85, 3.9, 0.75, CARD)
text(s, "similarity = cos(θ) = (A · B) / (|A| × |B|)",
     9.35, 2.0, 3.65, 0.45, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

score_data = [
    ("90–100%", "Excellent",  "Chunk directly answers the question",                  GREEN),
    ("70–89%",  "Good",       "Chunk is clearly relevant to the question",             BLUE),
    ("50–69%",  "Medium",     "Chunk is probably related — review the text",          YELLOW),
    ("30–49%",  "Weak",       "Loose connection — system struggled to find a match",  PINK),
    ("0–29%",   "No match",   "Unrelated — knowledge base may not cover this topic",  PURPLE),
]
text(s, "Similarity score guide:", 0.5, 3.1, 5, 0.35, size=13, bold=True, color=GREY)
for i, (score, label, desc, color) in enumerate(score_data):
    tpos = 3.5 + i * 0.62
    rect(s, 0.5, tpos, 1.15, 0.52, CARD)
    text(s, score, 0.55, tpos + 0.08, 1.08, 0.38,
         size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    rect(s, 1.73, tpos, 1.1, 0.52, CARD)
    text(s, label, 1.78, tpos + 0.08, 1.05, 0.38,
         size=11, color=color, align=PP_ALIGN.CENTER)
    text(s, desc, 2.95, tpos + 0.1, 6.5, 0.38, size=11, color=GREY)

row_of_cards(s, 3.1, [
    ("Cosine vs Euclidean",
     ["Cosine: measures direction (angle) — good for text.",
      "Euclidean: measures straight-line distance — biased",
      "  by vector magnitude (long docs look more distant).",
      "For text embeddings, cosine is almost always better.",
      "Set via: metadata={\"hnsw:space\": \"cosine\"}"],
     PURPLE),
    ("Dot Product",
     ["dot_product = sum(a[i] * b[i] for all i)",
      "For normalised vectors, equals cosine similarity.",
      "Faster to compute than cosine.",
      "Used by many vector DBs as default.",
      "Requires L2-normalised embeddings."],
     BLUE),
], card_h=2.15, gap=0.2)

rect(s, 9.2, 2.7, 3.9, 0.04, GREY)
for i, (label, val) in enumerate([
    ("Identical sentences:", "1.00  (cos 0°)"),
    ("Very similar meaning:", "0.90–0.99"),
    ("Related topic:", "0.65–0.89"),
    ("Unrelated:", "< 0.5  (cos > 60°)"),
]):
    rect(s, 9.2, 2.75 + i * 0.62, 3.9, 0.52, CARD)
    text(s, label, 9.35, 2.82 + i * 0.62, 2.0, 0.38, size=11, color=GREY)
    text(s, val, 11.4, 2.82 + i * 0.62, 1.6, 0.38,
         size=12, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)


# ── 7. VECTOR DATABASES ────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Vector Databases",
              "Storage and search engines built specifically for embeddings",
              bar_color=YELLOW)

dbs = [
    ("ChromaDB",
     ["Open-source, local-first", "Runs embedded in Python process",
      "No server needed", "Persisted as SQLite file",
      "Best for: local dev, small-medium scale",
      "Scale: up to ~1M vectors",
      "Language: Python native"],
     GREEN),
    ("FAISS",
     ["By Meta / Facebook AI Research", "In-memory (no persistence built-in)",
      "Extremely fast — C++ core", "Multiple index types (Flat, IVF, HNSW)",
      "Best for: research, large scale search",
      "Scale: 100M+ vectors",
      "Language: Python + C++"],
     BLUE),
    ("Pinecone",
     ["Fully managed cloud service", "No infrastructure to manage",
      "Real-time indexing", "Paid — usage-based pricing",
      "Best for: production cloud RAG apps",
      "Scale: billions of vectors",
      "Language: REST API + SDKs"],
     YELLOW),
    ("Weaviate",
     ["Open-source + cloud managed option", "Built-in vectorizer modules",
      "GraphQL query interface", "Hybrid search (vector + BM25)",
      "Best for: complex production RAG",
      "Scale: 10M–100M vectors",
      "Language: REST / Python / JS"],
     PURPLE),
    ("Qdrant",
     ["Open-source, written in Rust", "Extremely fast and memory efficient",
      "Rich filtering & payload support", "Docker or cloud deployment",
      "Best for: production with filtering needs",
      "Scale: 100M+ vectors",
      "Language: Python / REST / gRPC"],
     PINK),
    ("Milvus",
     ["Open-source, distributed", "Cloud-native, Kubernetes-ready",
      "Multiple index algorithms", "High availability built-in",
      "Best for: enterprise scale",
      "Scale: billions of vectors",
      "Language: Python / Java / Go"],
     GREEN),
]

for i, (name, bullets, color) in enumerate(dbs):
    col = i % 3
    row = i // 3
    lpos = 0.3 + col * 4.35
    tpos = 1.9  + row * 2.7
    card(s, lpos, tpos, 4.15, 2.5, name, bullets,
         title_color=color, body_size=10)


# ── 8. HNSW ────────────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "HNSW — Hierarchical Navigable Small World",
              "The algorithm that makes vector search fast — from O(n) to O(log n)",
              bar_color=PINK)

text(s, "The problem without HNSW:", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=YELLOW)
text(s,
     "Brute-force search compares every stored vector against the query vector. "
     "With 100,000 chunks of 768 dimensions, that means 76.8 million multiplications per query — "
     "taking seconds. HNSW reduces this to a few hundred comparisons.",
     0.5, 2.2, 12.3, 0.7, size=13, color=GREY)

row_of_cards(s, 3.0, [
    ("How HNSW builds the graph",
     ["Inserts vectors one by one into a layered graph.",
      "Layer 0: ALL vectors — fine-grained local connections.",
      "Layer 1: A subset — medium-range connections.",
      "Layer 2+: Fewer vectors — long-range 'highway' links.",
      "Like a highway system: motorway → A-road → local road.",
      "Built once at ingestion time. Searched on every query."],
     PINK),
    ("How HNSW searches",
     ["Start at a random entry node in the top layer.",
      "Greedily jump to whichever neighbour is closest to query.",
      "Drop down to the next layer and repeat.",
      "At Layer 0, explore the local neighbourhood precisely.",
      "Return the K nearest nodes found.",
      "Never checks the majority of stored vectors."],
     BLUE),
    ("Performance comparison",
     ["Brute force:  O(n × d)  — n=chunks, d=dimensions.",
      "HNSW:         O(log n × d).",
      "",
      "10,000 chunks:   ~4 vs 10,000 distance checks.",
      "1,000,000 chunks: ~7 vs 1,000,000 checks.",
      "Query time stays near-constant as DB grows."],
     GREEN),
    ("Key parameters",
     ["M: max connections per node (default 16).",
      "  Higher M = better recall, more memory.",
      "ef_construction: build quality (default 100).",
      "  Higher = slower build, better accuracy.",
      "ef: search quality (default 10).",
      "  Higher = slower search, more accurate."],
     YELLOW),
], card_h=2.9)

text(s,
     "ChromaDB uses HNSW internally by default. "
     "Set the distance metric with:  metadata={\"hnsw:space\": \"cosine\"}",
     0.5, 7.05, 12.3, 0.38, size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ── 9. TEXT CHUNKING STRATEGIES ────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Text Chunking Strategies",
              "How you split documents dramatically affects retrieval quality",
              bar_color=GREEN)

text(s, "Why chunking matters so much:", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=YELLOW)
text(s,
     "The embedding model converts each chunk to ONE vector. If a chunk contains too many "
     "topics, the vector becomes an average that matches nothing well. "
     "If chunks are too small, they lose context. Chunking is the biggest lever for RAG quality.",
     0.5, 2.2, 12.3, 0.7, size=13, color=GREY)

row_of_cards(s, 3.0, [
    ("Fixed-size Chunking",
     ["Split every N characters regardless of content.",
      "Simple and predictable.",
      "Overlap prevents sentence cutoffs.",
      "Best for: uniform documents, quick prototypes.",
      "Used in this project: chunk_size=600, overlap=100",
      "Risk: may cut mid-sentence or mid-concept."],
     BLUE),
    ("Recursive Character Splitting",
     ["Tries natural boundaries in priority order:",
      "  \\n\\n → \\n → . → space → character",
      "Falls back to next separator if chunk too large.",
      "Respects paragraph and sentence structure.",
      "Best for: general text, articles, reports.",
      "Used by LangChain's RecursiveCharacterTextSplitter."],
     GREEN),
    ("Semantic Chunking",
     ["Embeds every sentence first.",
      "Groups sentences whose embeddings are similar.",
      "Splits when meaning changes significantly.",
      "Produces topically coherent chunks.",
      "Best for: long documents with distinct sections.",
      "More expensive — requires embeddings upfront."],
     YELLOW),
    ("Document-Aware Chunking",
     ["Respects document structure: headings, sections.",
      "PDFs: split per page or per section heading.",
      "Code: split per function or class.",
      "HTML: split per <section> or <article> tag.",
      "Best for: structured content with clear hierarchy.",
      "Requires a structure-aware parser."],
     PINK),
], card_h=3.3)

text(s,
     "Overlap tip: 15–20% overlap is a good default. "
     "Too little = concepts split across chunks. Too much = duplicate retrieval.",
     0.5, 6.45, 12.3, 0.4, size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)

row_of_cards(s, 6.88, [
    ("Chunk too small  (<200 chars)",
     ["High precision but lacks context.",
      "LLM gets incomplete information.",
      "Use: very short factual Q&A"],
     GREY),
    ("Chunk ideal  (400–800 chars)",
     ["Balanced precision and context.",
      "Best for most document types.",
      "Start here and tune from results."],
     GREEN),
    ("Chunk too large  (>1500 chars)",
     ["Low search precision — too many topics.",
      "Wastes LLM context window.",
      "Use only with very focused documents."],
     GREY),
], card_h=0.98)


# ── 10. SHA-256 HASHING ────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Document Deduplication — SHA-256 Hashing",
              "Detecting changed files to avoid re-ingesting unchanged documents")

text(s, "What is a cryptographic hash?", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=BLUE)
text(s,
     "A hash function takes any input (a file, a string) and produces a fixed-length fingerprint. "
     "SHA-256 always produces a 64-character hex string. "
     "The tiniest change to the file — even one character — produces a completely different hash.",
     0.5, 2.2, 9.8, 0.75, size=13, color=GREY)

code_box(s, [
    "import hashlib",
    "sha = hashlib.sha256()",
    "with open('report.pdf', 'rb') as f:",
    "    for chunk in iter(lambda: f.read(8192), b''):",
    "        sha.update(chunk)",
    "fingerprint = sha.hexdigest()  # 64-char hex string",
], 10.4, 1.85, 2.7, 1.9, label="Python code")

row_of_cards(s, 3.1, [
    ("First ingestion",
     ["1. Compute hash of the file.",
      "2. Query ChromaDB: any chunks with this filename?",
      "3. No match found.",
      "4. Ingest all chunks, store hash in metadata.",
      "→ Status: ✓ added"],
     GREEN),
    ("Re-run, file unchanged",
     ["1. Compute hash of the file.",
      "2. Query ChromaDB: find chunks for this filename.",
      "3. Stored hash == current hash.",
      "4. Skip — no work to do.",
      "→ Status: – skipped  (saves time & API calls)"],
     BLUE),
    ("Re-run, file changed",
     ["1. Compute hash of the file.",
      "2. Query ChromaDB: find chunks for this filename.",
      "3. Stored hash != current hash.",
      "4. Delete old chunks, re-embed new version.",
      "→ Status: ↺ updated"],
     YELLOW),
    ("Why this matters",
     ["Embedding 1000 PDFs takes hours.",
      "Without deduplication: every run re-embeds all.",
      "With SHA-256: only changed files are re-embedded.",
      "Makes the pipeline safe to run repeatedly.",
      "→ Idempotent ingestion"],
     PINK),
], card_h=2.65)

text(s,
     "Other deduplication approaches: last-modified timestamp (faster but misses content-identical renames), "
     "MD5 hash (faster, slightly weaker collision resistance), content fingerprint on extracted text only.",
     0.5, 6.9, 12.3, 0.5, size=11, italic=True, color=GREY)


# ── 11. LLM TYPES FOR RAG ──────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "LLMs for RAG — Choosing the Right Model",
              "Cloud APIs vs local models — quality, cost, privacy, and latency trade-offs",
              bar_color=YELLOW)

llms = [
    ("OpenAI GPT-4o / GPT-4.1",
     ["Provider: OpenAI (cloud API)",
      "Context window: 128k tokens",
      "Quality: Best-in-class",
      "Cost: ~$2.50/$10 per 1M tokens",
      "Privacy: Data sent to OpenAI",
      "Best for: production apps where quality > cost"],
     BLUE),
    ("Anthropic Claude\n(claude-3.5-sonnet)",
     ["Provider: Anthropic (cloud API)",
      "Context window: 200k tokens",
      "Quality: Excellent reasoning + long context",
      "Cost: ~$3/$15 per 1M tokens",
      "Privacy: Data sent to Anthropic",
      "Best for: large-document RAG with long context"],
     GREEN),
    ("Google Gemini\n(gemini-1.5-pro)",
     ["Provider: Google (cloud API)",
      "Context window: 1 million tokens",
      "Quality: Excellent — largest context window",
      "Cost: ~$1.25/$5 per 1M tokens",
      "Privacy: Data sent to Google",
      "Best for: entire-document Q&A in one shot"],
     YELLOW),
    ("Meta Llama 3.x\n(local via Ollama)",
     ["Provider: Meta (open weights)",
      "Context window: 8k–128k tokens",
      "Quality: Good for its size class",
      "Cost: Free — runs on your hardware",
      "Privacy: 100% local — no data leaves machine",
      "Best for: private data, offline, zero cost"],
     PURPLE),
    ("Google Gemma 3\n(local via Ollama)",
     ["Provider: Google (open weights)",
      "Sizes: 1B / 4B / 12B / 27B parameters",
      "Quality: Good — efficient for size",
      "Cost: Free — runs on your hardware",
      "Privacy: 100% local",
      "Best for: low-RAM devices, quick prototypes"],
     PINK),
    ("Mistral / Mixtral\n(local or API)",
     ["Provider: Mistral AI",
      "Context window: 32k tokens",
      "Quality: Very efficient — punches above weight",
      "Cost: Free local, or low-cost API",
      "Privacy: Local option available",
      "Best for: multilingual RAG, balance of speed/quality"],
     GREEN),
]

for i, (name, bullets, color) in enumerate(llms):
    col = i % 3
    row = i // 3
    lpos = 0.3 + col * 4.35
    tpos = 1.85 + row * 2.55
    card(s, lpos, tpos, 4.15, 2.4, name, bullets,
         title_color=color, body_size=10)

text(s,
     "Key consideration: for private/sensitive documents, always use a local model (Llama, Gemma, Mistral). "
     "Cloud APIs process your data on external servers.",
     0.3, 7.1, 12.7, 0.38, size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── 12. SYSTEM PROMPT & GUARDRAILS ────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "System Prompts & Guardrails",
              "Preventing hallucination, confabulation, and off-topic answers",
              bar_color=PINK)

text(s, "How LLM message roles work:", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=PINK)

row_of_cards(s, 2.25, [
    ("system  (guardrail)",
     ["Sent before everything else.",
      "Sets rules, tone, and constraints.",
      "LLM treats this as standing instructions.",
      "Cannot be overridden by the user message."],
     PINK),
    ("user  (the question)",
     ["The actual question from the user.",
      "Contains the retrieved context blocks.",
      "This is the augmented prompt.",
      "Changes on every query."],
     BLUE),
    ("assistant  (the answer)",
     ["The LLM's previous responses.",
      "Used for multi-turn conversations.",
      "In single-turn RAG this is empty.",
      "Enables follow-up questions on context."],
     GREEN),
], card_h=1.65)

text(s, "A well-designed RAG system prompt:", 0.5, 4.1, 7, 0.35,
     size=13, bold=True, color=PINK)
code_box(s, [
    "You are a helpful assistant that answers questions",
    "strictly based on the provided context.",
    "",
    "Rules:",
    "- Only use information from the context below to answer.",
    "- If the context does not contain enough information,",
    '  say "I don\'t have enough information in the provided',
    '  documents to answer this."',
    "- Be concise and accurate.",
    "- Cite the source document name when relevant.",
], 0.4, 4.5, 7.5, 2.85)

row_of_cards(s, 4.1, [
    ("Anti-hallucination",
     ["Without: LLM answers from training data.",
      "With: LLM only uses provided context.",
      "Result: answers are always verifiable",
      "against your actual documents."],
     YELLOW),
    ("Anti-confabulation",
     ["Without: LLM guesses confidently when unsure.",
      "With: LLM says 'I don't know' explicitly.",
      "Result: users know when to look elsewhere.",
      "Trust is maintained."],
     BLUE),
    ("Traceability",
     ["Forces the model to cite source files.",
      "User can verify: 'According to report.pdf...'",
      "Critical for legal, medical, compliance use.",
      "Auditable AI — know where every fact came from."],
     GREEN),
], card_h=2.35, gap=0.15)


# ── 13. AUGMENTED PROMPT ──────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "The Augmented Prompt",
              "Exactly what the LLM reads — the core of the RAG pattern")

text(s, "Structure of the augmented prompt:", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=BLUE)

code_box(s, [
    "Context:",
    "",
    "[Source 1: annual_report.pdf, page 3, similarity 87.2%]",
    "...text of the most relevant chunk (up to chunk_size characters)...",
    "",
    "---",
    "",
    "[Source 2: q3_results.pdf, page 1, similarity 74.1%]",
    "...text of the second most relevant chunk...",
    "",
    "---",
    "",
    "[Source 3: notes.txt, page 0, similarity 61.3%]",
    "...text of the third most relevant chunk...",
    "",
    "Question: What was the Q3 revenue and how did it compare to Q2?",
], 0.4, 2.3, 7.8, 4.9)

row_of_cards(s, 2.3, [
    ("Why label each source?",
     ["LLM sees which file each chunk came from.",
      "Enables it to cite sources in the answer.",
      "Helps debug poor answers — which source failed?",
      "Similarity % tells LLM how confident to be."],
     BLUE),
    ("Top-K trade-off",
     ["K=2: Fast, focused, may miss relevant info.",
      "K=4: Good default for most use cases.",
      "K=8: More context, slower, may dilute answer.",
      "K too high: LLM gets confused by noise.",
      "Tune K based on document type and query complexity."],
     GREEN),
    ("Token streaming",
     ["LLM generates one token (~1 word) at a time.",
      "stream=True returns tokens as they're generated.",
      "UI displays each word as it arrives.",
      "Faster perceived response time.",
      "Blinking ▌ cursor shows generation in progress."],
     YELLOW),
], card_h=4.9, gap=0.2)


# ── 14. RAG VARIANTS ──────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "RAG Variants",
              "From basic to advanced architectures — choosing the right approach",
              bar_color=GREEN)

row_of_cards(s, 1.85, [
    ("Naive RAG\n(Basic)",
     ["The pattern covered in this POC.",
      "Ingest → Embed → Store → Retrieve → Generate.",
      "Single-step retrieval, no refinement.",
      "Good for: prototypes, simple document Q&A.",
      "Weakness: no query understanding, fixed K.",
      "Weakness: no feedback loop or re-ranking."],
     BLUE),
    ("Advanced RAG",
     ["Adds pre-retrieval and post-retrieval steps.",
      "Query rewriting: rephrase question for better search.",
      "Re-ranking: re-score retrieved chunks with a cross-encoder.",
      "HyDE: generate a hypothetical answer, embed that.",
      "Good for: production apps needing higher accuracy.",
      "More complex but significantly better results."],
     GREEN),
    ("Modular RAG",
     ["Each component is independently swappable.",
      "Mix any embedding model + any vector DB + any LLM.",
      "Add modules: routing, fusion, self-reflection.",
      "Query routing: send to different pipelines based on topic.",
      "Good for: enterprise systems with multiple document types.",
      "Frameworks: LangChain, LlamaIndex, Haystack."],
     YELLOW),
    ("Graph RAG\n(Microsoft)",
     ["Builds a knowledge graph from documents.",
      "Extracts entities and relationships (not just chunks).",
      "Can answer multi-hop questions across documents.",
      "Example: 'How does X relate to Y in the annual report?'",
      "Good for: complex relational queries across large corpora.",
      "Much more expensive to build and query."],
     PINK),
], card_h=2.8)

row_of_cards(s, 4.8, [
    ("Agentic RAG",
     ["LLM decides when and what to retrieve.",
      "Can call retrieval multiple times per question.",
      "Tool-use: search, calculate, look up APIs.",
      "Self-corrects if initial answer seems incomplete.",
      "Best for: complex, multi-step research questions."],
     PURPLE),
    ("Hybrid Search RAG",
     ["Combines vector search + BM25 keyword search.",
      "BM25 is great for exact terms (names, product IDs).",
      "Vector search is great for semantic meaning.",
      "Fusion: re-rank results from both methods.",
      "Best for: documents with specific jargon or IDs."],
     BLUE),
    ("Multimodal RAG",
     ["Handles images, tables, charts alongside text.",
      "Uses vision-language models (GPT-4o, LLaVA).",
      "Embeds images as vectors using CLIP or similar.",
      "Retrieves both text and image chunks.",
      "Best for: product catalogues, scientific papers."],
     GREEN),
    ("Self-RAG",
     ["LLM reflects on its own outputs.",
      "Generates 'retrieval tokens' to decide when to retrieve.",
      "Critiques its answer with 'critique tokens'.",
      "Retries retrieval if answer quality is low.",
      "Best for: applications requiring very high accuracy."],
     YELLOW),
], card_h=2.55)


# ── 15. RE-RANKING ────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Re-ranking",
              "Improving retrieval precision with a second-pass scoring model",
              bar_color=PURPLE)

text(s, "Why retrieval alone isn't always enough:", 0.5, 1.85, 12, 0.35,
     size=14, bold=True, color=YELLOW)
text(s,
     "Bi-encoder embedding search (what we use) embeds query and document independently — it's fast "
     "but misses fine-grained interactions between query words and document words. "
     "A cross-encoder re-ranker reads BOTH query and chunk together, scoring their relevance far more accurately.",
     0.5, 2.2, 12.3, 0.75, size=13, color=GREY)

text(s, "Two-stage pipeline:", 0.5, 3.05, 12, 0.35, size=14, bold=True, color=PURPLE)

stages = [
    ("Stage 1 — Recall\n(Bi-encoder / Vector Search)",
     ["Fast: embeds query and docs separately.",
      "Retrieves Top-50 or Top-100 candidates.",
      "Uses approximate nearest neighbour (HNSW).",
      "Takes milliseconds.",
      "May include some loosely relevant results.",
      "Model: nomic-embed-text, OpenAI embeddings, etc."],
     BLUE, 0.4),
    ("Stage 2 — Precision\n(Cross-encoder / Re-ranker)",
     ["Slow but accurate: reads query + chunk together.",
      "Scores all 50–100 candidates with joint attention.",
      "Re-orders by true relevance.",
      "Takes ~200ms–1s.",
      "Returns Top-K truly relevant results.",
      "Model: cross-encoder/ms-marco-MiniLM-L-6-v2"],
     PURPLE, 4.7),
    ("Result — LLM sees only\nhigh-quality context",
     ["Only the re-ranked Top-K go into the prompt.",
      "Significantly reduces hallucination.",
      "Improves answer faithfulness.",
      "Worth the extra latency for production apps.",
      "Can skip in simple/fast scenarios.",
      "Tools: Cohere Rerank API, sentence-transformers"],
     GREEN, 9.0),
]
for title, bullets, color, lpos in stages:
    card(s, lpos, 3.5, 4.1, 3.2, title, bullets, title_color=color)
for x in [4.55, 8.85]:
    text(s, "→", x, 4.75, 0.2, 0.5,
         size=28, bold=True, color=GREY, align=PP_ALIGN.CENTER)

row_of_cards(s, 6.85, [
    ("When to add re-ranking",
     ["Retrieval quality is mediocre (scores 50–70%).",
      "Users report irrelevant answers.",
      "High-stakes domain: legal, medical, finance.",
      "Documents have overlapping terminology."],
     YELLOW),
    ("Popular re-ranker models",
     ["cross-encoder/ms-marco-MiniLM-L-6-v2  (fast, local)",
      "BAAI/bge-reranker-v2-m3  (multilingual)",
      "Cohere Rerank 3  (best quality, cloud API)",
      "mixedbread-ai/mxbai-rerank-large  (open-source)"],
     PURPLE),
], card_h=0.9)


# ── 16. RAG EVALUATION METRICS ────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Evaluating RAG — Key Metrics",
              "How to measure if your RAG system is actually working well",
              bar_color=GREEN)

text(s, "The RAG evaluation triangle — 3 components to measure:",
     0.5, 1.85, 12, 0.35, size=14, bold=True, color=GREEN)

row_of_cards(s, 2.25, [
    ("Retrieval Quality",
     ["Does the retriever find the right chunks?",
      "",
      "Context Recall: were relevant chunks retrieved?",
      "Context Precision: were retrieved chunks all relevant?",
      "  (no irrelevant chunks in Top-K?)",
      "Hit Rate: was the correct chunk in Top-K at all?",
      "MRR: was the best chunk ranked #1?"],
     BLUE),
    ("Answer Quality",
     ["Is the generated answer good?",
      "",
      "Faithfulness: does the answer stay within the context?",
      "  (measures hallucination rate)",
      "Answer Relevancy: does the answer actually address",
      "  the question asked?",
      "Answer Correctness: is it factually right?"],
     GREEN),
    ("End-to-End Quality",
     ["Does the whole pipeline work together?",
      "",
      "RAGAS score: weighted average of above metrics.",
      "Human evaluation: domain expert review.",
      "A/B testing: compare two RAG configurations.",
      "LLM-as-judge: use GPT-4 to score answers",
      "  against reference answers."],
     YELLOW),
], card_h=3.2)

row_of_cards(s, 5.6, [
    ("Faithfulness (most important)",
     ["Measures: does the answer contain ONLY information",
      "from the retrieved context?",
      "Score 0–1. High score = no hallucination.",
      "Tool: RAGAS library",
      "ragas.metrics.faithfulness"],
     PINK),
    ("Context Precision",
     ["Measures: are retrieved chunks relevant to the query?",
      "High precision = no noise in the Top-K results.",
      "Low precision = re-ranking or better chunking needed.",
      "Tool: RAGAS library",
      "ragas.metrics.context_precision"],
     PURPLE),
    ("Answer Relevancy",
     ["Measures: does the answer address the question?",
      "High = directly answers what was asked.",
      "Low = answer is technically correct but off-topic.",
      "Tool: RAGAS library",
      "ragas.metrics.answer_relevancy"],
     BLUE),
    ("Tooling",
     ["RAGAS: pip install ragas",
      "  — automated metrics using LLM as judge",
      "DeepEval: alternative evaluation framework",
      "TruLens: evaluation + observability",
      "LangSmith: LangChain's tracing + eval platform"],
     GREEN),
], card_h=1.72)


# ── 17. PRODUCTION CONSIDERATIONS ─────────────────────────────────────────────
s = blank(prs)
set_bg(s)
slide_heading(s, "Production RAG — Best Practices",
              "Moving from prototype to a reliable, scalable system",
              bar_color=BLUE)

row_of_cards(s, 1.85, [
    ("Metadata Filtering",
     ["Add metadata tags to chunks at ingestion time.",
      "Filter at query time: only search within dept='finance'",
      "  or date > '2024-01-01'.",
      "Dramatically improves precision in large corpora.",
      "ChromaDB: collection.query(where={'dept': 'finance'})",
      "Pinecone: filter={'dept': {'$eq': 'finance'}}"],
     BLUE),
    ("Hybrid Search",
     ["Combine dense (vector) + sparse (BM25) retrieval.",
      "BM25 excels at exact matches: names, product codes, IDs.",
      "Vector excels at semantic meaning.",
      "Reciprocal Rank Fusion (RRF) merges both rankings.",
      "Implemented in Weaviate, Qdrant, Elasticsearch.",
      "Recommended for production systems."],
     GREEN),
    ("Query Rewriting",
     ["Rephrase vague or short queries before embedding.",
      "Example: 'revenue?' → 'What was the annual revenue?'",
      "Use LLM to expand abbreviations and add context.",
      "Multi-query: generate 3–5 query variants, merge results.",
      "HyDE: embed a hypothetical answer to the question.",
      "All improve recall for ambiguous questions."],
     YELLOW),
    ("Observability",
     ["Log every query: question + retrieved chunks + answer.",
      "Track latency per step: embed / search / generate.",
      "Monitor faithfulness and relevancy over time.",
      "Alert on low-similarity retrievals (< 40%).",
      "Tools: LangSmith, Helicone, Arize Phoenix.",
      "Essential for debugging production issues."],
     PINK),
], card_h=2.85)

row_of_cards(s, 4.85, [
    ("Chunking tuning",
     ["Start: chunk_size=600, overlap=100.",
      "Short factual docs: reduce to 300–400.",
      "Long narrative docs: increase to 800–1000.",
      "Evaluate with RAGAS after each change.",
      "Re-ingest after every chunking change."],
     PURPLE),
    ("Embedding model upgrade",
     ["Start local (nomic-embed-text) for development.",
      "Upgrade to mxbai-embed-large for better local quality.",
      "Use OpenAI text-embedding-3-large for max accuracy.",
      "Always re-ingest EVERYTHING when switching models.",
      "Benchmark on your own data before committing."],
     BLUE),
    ("Security",
     ["Never store PII in chunk metadata.",
      "Use access control on the vector DB.",
      "For sensitive data: local LLM + local embedding only.",
      "Sanitise retrieved chunks before sending to cloud LLMs.",
      "Audit logs for compliance (who asked what)."],
     PINK),
    ("Scaling",
     ["< 100k chunks: ChromaDB or FAISS — works great.",
      "100k–10M chunks: Qdrant or Weaviate self-hosted.",
      "> 10M chunks: Pinecone or Milvus cloud.",
      "Cache embeddings of frequent queries.",
      "Pre-warm the HNSW index on startup."],
     GREEN),
], card_h=2.5)


# ── 18. SUMMARY ───────────────────────────────────────────────────────────────
s = blank(prs)
set_bg(s)
top_bar(s, BLUE)
text(s, "Everything Covered", 0.5, 0.2, 12, 0.75,
     size=30, bold=True, color=WHITE)

all_topics = [
    ("RAG",                "Retrieve → Augment → Generate. Grounds LLM answers in your documents.",                                              BLUE),
    ("Embeddings",         "Text → fixed-length number vectors. Similar meaning = similar vectors.",                                              GREEN),
    ("Embedding Models",   "OpenAI, Cohere, Google, HuggingFace (BGE/E5), Ollama (nomic, mxbai).",                                               YELLOW),
    ("Cosine Similarity",  "Angle between vectors. 0 = unrelated, 1 = identical meaning.",                                                       PURPLE),
    ("Vector Databases",   "ChromaDB, FAISS, Pinecone, Weaviate, Qdrant, Milvus — each with different scale/trade-offs.",                       BLUE),
    ("HNSW",               "Hierarchical graph algorithm. O(log n) search instead of O(n). Fast nearest-neighbour lookup.",                       PINK),
    ("Text Chunking",      "Fixed-size, recursive, semantic, document-aware strategies. Size + overlap are critical.",                            GREEN),
    ("SHA-256 Hashing",    "File fingerprinting for deduplication. Only re-embed changed files.",                                                 YELLOW),
    ("LLMs for RAG",       "GPT-4o, Claude, Gemini (cloud). Llama, Gemma, Mistral (local). Cost vs quality vs privacy.",                        BLUE),
    ("Augmented Prompt",   "Context blocks + source labels + question = exactly what the LLM reads.",                                            GREEN),
    ("Guardrails",         "System prompt rules: anti-hallucination, anti-confabulation, traceability.",                                         PINK),
    ("Token Streaming",    "LLM streams tokens word-by-word for faster perceived response.",                                                      PURPLE),
    ("RAG Variants",       "Naive → Advanced → Modular → Graph → Agentic → Hybrid → Multimodal → Self-RAG.",                                    YELLOW),
    ("Re-ranking",         "Two-stage retrieval: fast bi-encoder recall + slow cross-encoder precision.",                                         PURPLE),
    ("Evaluation Metrics", "Faithfulness, context precision/recall, answer relevancy, RAGAS score.",                                              GREEN),
    ("Production Best Practices", "Metadata filtering, hybrid search, query rewriting, observability, security.",                                BLUE),
]

cols = 2
per_col = len(all_topics) // cols + 1
for i, (name, desc, color) in enumerate(all_topics):
    col = i // per_col
    row = i % per_col
    lpos = 0.3  + col * 6.6
    tpos = 1.05 + row * 0.76
    rect(s, lpos, tpos, 6.45, 0.68, CARD)
    text(s, name, lpos + 0.12, tpos + 0.06, 2.1, 0.56,
         size=12, bold=True, color=color)
    text(s, desc, lpos + 2.3, tpos + 0.1, 4.0, 0.52,
         size=10, color=GREY)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/paritosh/development/RAG- POC/RAG_POC_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
print(f"Total slides: {len(prs.slides)}")
